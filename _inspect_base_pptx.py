"""Inspect the base PPTX to extract slide content."""
from pptx import Presentation

prs = Presentation("paper/SHARP-LLM_A_Framework_for_Vulnerabilitty_Detection_AM.SC.R4CSE25007-PT.pptx")
for i, slide in enumerate(prs.slides):
    layout_name = slide.slide_layout.name
    print(f"--- Slide {i+1} ({layout_name}) ---")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    print(f"  {t[:180]}")
        elif shape.has_table:
            tbl = shape.table
            nr = len(list(tbl.rows))
            nc = len(tbl.columns)
            for r in range(nr):
                row_cells = []
                for c in range(nc):
                    row_cells.append(tbl.cell(r, c).text.strip()[:35])
                print("  | " + " | ".join(row_cells) + " |")
    print()
