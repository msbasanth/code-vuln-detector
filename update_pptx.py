"""
Update SHARP-LLM PowerPoint presentation.
Keeps template (slide 1: cover, slide 2: title) and rebuilds content slides
with sections: Introduction, Literature Review & Research Gap, Novelty/Problem Statement,
Methodology, Block Diagram, Results, References.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy
from lxml import etree

SRC = r'd:\Repositories\code-vuln-detector\paper\SHARP-LLM_Secure_Lightweight_Framework_for_Source_Code_Vulnerability_Detection_and_Risk_Prioritization_in_Healthcare_Software_AM.SC.R4CSE25007-PT_v1.0.pptx'
DST = r'd:\Repositories\code-vuln-detector\paper\SHARP-LLM_Secure_Lightweight_Framework_for_Source_Code_Vulnerability_Detection_and_Risk_Prioritization_in_Healthcare_Software_AM.SC.R4CSE25007-PT_v1.0.pptx'

prs = Presentation(SRC)

# ── Helpers ──────────────────────────────────────────────────────────────

ACCENT = RGBColor(0xA4, 0x12, 0x3F)   # section-number colour
DARK   = RGBColor(0x00, 0x00, 0x00)   # body text
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def _delete_slides_from(prs, keep_count):
    """Delete all slides after the first *keep_count* slides."""
    sldIdLst = prs._element.sldIdLst
    ids = list(sldIdLst)
    for sldId in ids[keep_count:]:
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def _clone_slide(prs, template_slide):
    """Deep-clone *template_slide* and append it to *prs*."""
    # Copy the XML of the slide
    slide_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Remove all shapes from the new slide
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Deep copy every shape from the template
    for shape in template_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    return new_slide


def add_section_divider(prs, number, title_text, template_divider):
    """Add a section divider slide cloned from an existing divider template."""
    slide = _clone_slide(prs, template_divider)
    # Update the number
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                full = para.text.strip()
                if full in ('01', '02', '03', '04', '05', '06', '07', '08'):
                    # Replace number
                    for run in para.runs:
                        run.text = ''
                    if para.runs:
                        para.runs[0].text = number
                    break
                # check for section title pattern (bold, ~32pt)
            for para in shape.text_frame.paragraphs:
                full = para.text.strip()
                for run in para.runs:
                    if run.font.bold and run.font.size and run.font.size >= Pt(28):
                        # This is the section title
                        pass
    # Now set the number text and title text
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        paras = list(shape.text_frame.paragraphs)
        for para in paras:
            txt = ''.join(r.text for r in para.runs).strip()
            # Detect the number text box (single big number)
            if txt in ('01', '02', '03', '04', '05', '06', '07', '08'):
                for run in para.runs:
                    run.text = ''
                para.runs[0].text = number
            # Detect the title text box (bold, larger, descriptive)
            for run in para.runs:
                if run.font.bold and run.font.size and run.font.size >= Pt(28) and len(txt) > 3:
                    # This is the section title
                    for r in para.runs:
                        r.text = ''
                    para.runs[0].text = title_text
                    break
    return slide


def add_content_slide(prs, title_text, bullets, template_content):
    """Add a content slide cloned from template_content. Replaces title and body."""
    slide = _clone_slide(prs, template_content)

    # Find and update title placeholder
    for shape in slide.shapes:
        if shape.shape_type == 14:  # PLACEHOLDER
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ''
                    if para.runs:
                        para.runs[0].text = title_text
                    else:
                        run = para.add_run()
                        run.text = title_text
                # Only update first para, clear the rest
                paras = list(shape.text_frame.paragraphs)
                for p in paras[1:]:
                    pp = p._p
                    pp.getparent().remove(pp)
                break

    # Find the main content TextBox (the largest one, not the title)
    content_box = None
    for shape in slide.shapes:
        if shape.shape_type == 17 and shape.has_text_frame:  # TEXT_BOX
            if shape.width > Emu(8000000):  # The wide content box
                content_box = shape
                break

    if content_box:
        tf = content_box.text_frame
        # Clear existing paragraphs
        for para in list(tf.paragraphs):
            pp = para._p
            pp.getparent().remove(pp)

        # Get reference styling from first paragraph if exists
        for i, bullet in enumerate(bullets):
            p = etree.SubElement(tf._txBody, '{http://schemas.openxmlformats.org/drawingml/2006/main}p')

            # Check if it's a sub-bullet
            is_sub = bullet.startswith('     ')
            text = bullet.strip()

            r = etree.SubElement(p, '{http://schemas.openxmlformats.org/drawingml/2006/main}r')
            rPr = etree.SubElement(r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr', attrib={
                'lang': 'en-US',
                'dirty': '0',
            })
            if is_sub:
                rPr.set('sz', '2159')  # 17pt for sub-bullets
            else:
                rPr.set('sz', '2540')  # 20pt for main bullets

            # Set font
            latin = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}latin',
                                     attrib={'typeface': 'Calibri'})

            t = etree.SubElement(r, '{http://schemas.openxmlformats.org/drawingml/2006/main}t')
            t.text = text

    return slide


def add_table_slide(prs, title_text, headers, rows, template_content):
    """Add a slide with a table. Clone a content slide, remove the content box, add a table."""
    slide = _clone_slide(prs, template_content)

    # Update title
    for shape in slide.shapes:
        if shape.shape_type == 14:  # PLACEHOLDER
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ''
                    if para.runs:
                        para.runs[0].text = title_text
                    else:
                        run = para.add_run()
                        run.text = title_text
                paras = list(shape.text_frame.paragraphs)
                for p in paras[1:]:
                    pp = p._p
                    pp.getparent().remove(pp)
                break

    # Remove the main content TextBox
    for shape in list(slide.shapes):
        if shape.shape_type == 17 and shape.has_text_frame:
            if shape.width > Emu(8000000):
                sp = shape._element
                sp.getparent().remove(sp)
                break

    # Add table
    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)
    left = Emu(1295400)
    top = Emu(1810978)
    width = Emu(9601200)
    height = Emu(min(4261659, 500000 * num_rows))

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # Style header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.bold = True
            para.font.name = 'Calibri'
            para.font.color.rgb = WHITE
        # Header background
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = ACCENT

    # Fill data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(12)
                para.font.name = 'Calibri'
                para.font.bold = False

    return slide


# ── Capture template slides before deletion ─────────────────────────────
# Slide 3 (index 2) = section divider template
# Slide 4 (index 3) = content slide template
divider_template_xml = copy.deepcopy(prs.slides[2]._element)
content_template_xml = copy.deepcopy(prs.slides[3]._element)

# We need the actual slide objects, so save references
divider_slide = prs.slides[2]   # Section divider (01 - Context of the Problem)
content_slide = prs.slides[3]   # Content (Why Vulnerability Detection Matters)
table_slide_ref = prs.slides[14] # Table slide (Model Comparison)

# ── Delete all slides after the first 2 (keep cover + title) ────────────
_delete_slides_from(prs, 2)

# ══════════════════════════════════════════════════════════════════════════
#  SECTION 1: INTRODUCTION
# ══════════════════════════════════════════════════════════════════════════
add_section_divider(prs, '01', 'Introduction', divider_slide)

add_content_slide(prs, 'Why Vulnerability Detection Matters', [
    '▸  Software vulnerabilities cause security breaches, service disruption, and data exposure',
    '▸  Healthcare software is especially critical \u2014 vulnerabilities can compromise:',
    '     ◦  Electronic Health Records (EHR) and patient data confidentiality',
    '     ◦  Clinical decision support and diagnostic workflows',
    '     ◦  Telemedicine platforms and connected medical devices',
    '     ◦  Regulatory compliance (HIPAA, GDPR)',
    '▸  Traditional SAST tools have high false-positive rates and limited generalization',
    '▸  LLM-based approaches show promise but are often computationally expensive',
    '▸  No existing framework combines lightweight detection with healthcare-oriented risk prioritization',
], content_slide)

add_content_slide(prs, 'What is SHARP-LLM?', [
    '▸  Secure Healthcare-Aligned Risk Prediction via Lightweight LLMs',
    '▸  A secure and lightweight framework for source code vulnerability detection',
    '     ◦  Classifies C/C++ code into 118 CWE categories',
    '     ◦  Fine-tuned lightweight LLMs (35M\u2013125M parameters)',
    '     ◦  Runs on consumer-grade GPU (8 GB VRAM)',
    '▸  Healthcare-oriented risk prioritization layer:',
    '     ◦  Maps CWEs to LINDDUN privacy threat types',
    '     ◦  Links to policy-relevant control domains',
    '     ◦  Produces ordinal priority levels (Critical / High / Medium / Low)',
    '▸  Sub-30 ms inference latency \u2014 viable for CI/CD, IDE plugins, real-time code review',
], content_slide)

# ══════════════════════════════════════════════════════════════════════════
#  SECTION 2: LITERATURE REVIEW & RESEARCH GAP
# ══════════════════════════════════════════════════════════════════════════
add_section_divider(prs, '02', 'Literature Review & Research Gap', divider_slide)

add_table_slide(prs, 'Related Work \u2014 Existing Approaches', [
    'Framework', 'Authors', 'Focus', 'Limitation',
], [
    ['GRACE', 'Lu et al. (2024)', 'Graph-enhanced LLM + ICL', 'Prompt sensitivity, limited languages'],
    ['DLAP', 'Yang et al. (2025)', 'DL-guided prompting', 'Dependent on DL model quality'],
    ['FUSEVUL', 'Tian et al. (2025)', 'Multi-modal detection', 'Limited to C/C++, basic prompting'],
    ['LLMVulExp', 'Mao et al. (2025)', 'Explainable via teacher-student', 'High annotation & compute cost'],
    ['VulACLLM', 'Liu et al. (2025)', 'AST+CFG+LoRA/distillation', 'Performance trade-off from compression'],
    ['SecureFalcon', 'Ferrag et al. (2025)', 'Compact LLM classification', 'Limited to known patterns'],
    ['RLV', 'Qiu et al. (2026)', 'Repository-level context', 'Context-length limits'],
    ['GPTVD', 'Chen et al. (2026)', 'Static slicing + CoT', 'Limited dataset, scalability'],
], content_slide)

add_content_slide(prs, 'Three Key Research Gaps', [
    '▸  Gap 1: Computational Cost',
    '     ◦  Many LLM-based methods rely on expensive architectures (large GPUs, >10B params)',
    '     ◦  Limits practical deployment in routine secure development workflows',
    '',
    '▸  Gap 2: No Domain-Specific Prioritization',
    '     ◦  Literature focuses on detection accuracy alone',
    '     ◦  No framework addresses how findings should be prioritized in healthcare settings',
    '     ◦  A buffer overflow in an EHR system \u2260 a buffer overflow in a game engine',
    '',
    '▸  Gap 3: Limited Healthcare-Oriented Resources',
    '     ◦  Privacy-aware vulnerability resources for healthcare are scarce',
    '     ◦  Missing link: translating technical CWE findings into healthcare-relevant risk signals',
], content_slide)

# ══════════════════════════════════════════════════════════════════════════
#  SECTION 3: NOVELTY / PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════════════
add_section_divider(prs, '03', 'Novelty / Problem Statement', divider_slide)

add_content_slide(prs, 'Problem Statement', [
    '▸  Source code vulnerability detection in C/C++ programs using fine-tuned lightweight LLMs',
    '',
    '▸  Given a code sample, determine whether it contains vulnerability-relevant patterns',
    '     ◦  Classify into one of 118 CWE categories',
    '',
    '▸  Additionally, prioritize findings in healthcare contexts where defects may impact:',
    '     ◦  Patient data confidentiality and integrity',
    '     ◦  Service continuity and system trust',
    '     ◦  Regulatory compliance',
    '',
    '▸  Three key challenges:',
    '     ◦  (a) Vulnerable vs. safe code may differ in only minor semantic details',
    '     ◦  (b) 118 CWE categories with uneven distributions (1,500 vs. 5 samples)',
    '     ◦  (c) High-performing models must be practical on consumer hardware',
], content_slide)

add_content_slide(prs, 'Three-Fold Novelty', [
    '▸  Novelty 1: Template-Aware Data Splitting',
    '     ◦  Eliminates structural data leakage inherent in the Juliet Test Suite',
    '     ◦  Splits at template level using GroupShuffleSplit \u2014 zero template overlap',
    '     ◦  First rigorous anti-leakage split for Juliet-based evaluations',
    '',
    '▸  Novelty 2: Multi-Paradigm Model Evaluation',
    '     ◦  First study to compare 3 paradigms on 118-class CWE classification:',
    '     ◦  Encoder fine-tuning | QLoRA fine-tuning | Zero-shot generative inference',
    '',
    '▸  Novelty 3: Healthcare-Oriented Risk Prioritization Layer',
    '     ◦  Three-stage post-detection pipeline:',
    '     ◦  CWE \u2192 LINDDUN threat types \u2192 Policy control domains \u2192 Weighted risk score',
    '     ◦  P = \u03bb\u2081s + \u03bb\u2082E + \u03bb\u2083D + \u03bb\u2084S + \u03bb\u2085C\u2019 \u2192 Priority Level',
], content_slide)

# ══════════════════════════════════════════════════════════════════════════
#  SECTION 4: METHODOLOGY
# ══════════════════════════════════════════════════════════════════════════
add_section_divider(prs, '04', 'Methodology', divider_slide)

add_content_slide(prs, 'Dataset & Preprocessing', [
    '▸  NIST Juliet Test Suite for C/C++ v1.3',
    '     ◦  105,183 source code samples across 118 CWE categories',
    '     ◦  Training: 82,736 samples (1,323 templates)',
    '     ◦  Test: 22,447 samples (366 templates, 87 CWEs)',
    '     ◦  Zero template overlap (80/20 split, seed=42)',
    '',
    '▸  4-Stage Preprocessing Pipeline:',
    '     ◦  Stage 1: Header removal (template-generated comment blocks)',
    '     ◦  Stage 2: Comment stripping (state machine, preserves string literals)',
    '     ◦  Stage 3: Guard removal (#ifndef OMITBAD, #ifndef OMITGOOD, etc.)',
    '     ◦  Stage 4: Whitespace normalization (compact representation)',
], content_slide)

add_table_slide(prs, 'Models Evaluated', [
    'Model', 'Paradigm', 'Params', 'Size (MB)', 'Key Details',
], [
    ['CodeT5-Small', 'Encoder FT', '35.4M', '405', 'Full fine-tuning, FP16'],
    ['CodeT5-Base', 'Encoder FT', '109.7M', '1,256', 'Full fine-tuning, FP16'],
    ['CodeBERT-Base', 'Encoder FT', '124.7M', '1,423', 'Full fine-tuning, FP16'],
    ['GraphCodeBERT', 'Encoder FT', '124.7M', '1,423', 'Full fine-tuning, FP16'],
    ['CodeGemma-2B', 'QLoRA FT', '2.51B', '4,677', '4-bit NF4, LoRA r=16, \u03b1=32'],
    ['CodeGemma-2B', 'Zero-Shot', '2.51B', '4,677', 'Beam search B=10'],
    ['Gemma-4-E2B-IT', 'Zero-Shot', '5.13B', '9,561', 'Instruction-tuned, B=2'],
], content_slide)

add_content_slide(prs, 'Training Configuration', [
    '▸  Classification Head: Mean pooling \u2192 Dropout \u2192 Linear \u2192 Softmax \u2192 118-dim',
    '',
    '▸  Encoder Models:',
    '     ◦  Learning rate: 5\u00d710\u207b\u2075, AdamW optimizer, batch size 8',
    '     ◦  Max sequence length: 256 tokens, 2 epochs, FP16 precision',
    '     ◦  Early stopping: patience=1',
    '',
    '▸  QLoRA (CodeGemma-2B):',
    '     ◦  Learning rate: 2\u00d710\u207b\u2074, PagedAdamW 8-bit, effective batch 8',
    '     ◦  LoRA targets: q, k, v, o, gate, up, down projections',
    '     ◦  Memory: ~5.5 GB (FP16) \u2192 ~1.8 GB (NF4), only ~0.5% params trained',
    '',
    '▸  Hardware: NVIDIA RTX 2000 Ada (8 GB VRAM), Intel i7-13850HX',
    '▸  Stack: Python 3.13, PyTorch \u22652.0, HuggingFace Transformers \u22655.5',
], content_slide)

add_content_slide(prs, 'Healthcare Risk Prioritization Layer', [
    '▸  Step 1 \u2014 Privacy & Security Risk Categorization:',
    '     ◦  R = f(c) \u2014 map detected CWE to LINDDUN threat types',
    '     ◦  Linkability, Identifiability, Non-repudiation, Detectability,',
    '     ◦  Information Disclosure, Unawareness, Non-compliance',
    '',
    '▸  Step 2 \u2014 Policy-Relevant Control Domains:',
    '     ◦  C = g(R) \u2014 confidentiality, access control, audit & accountability,',
    '     ◦  secure data handling, service availability',
    '',
    '▸  Step 3 \u2014 Weighted Risk Prioritization:',
    '     ◦  P = w\u2081\u00b7s + w\u2082\u00b7E + w\u2083\u00b7D + w\u2084\u00b7S + w\u2085\u00b7C\u2019',
    '     ◦  s = model confidence, E = severity, D = data sensitivity',
    '     ◦  S = service/safety impact, C\u2019 = control domain importance',
    '',
    '▸  Output: Ordinal priority level (Critical / High / Medium / Low)',
], content_slide)

# ══════════════════════════════════════════════════════════════════════════
#  SECTION 5: BLOCK DIAGRAM
# ══════════════════════════════════════════════════════════════════════════
add_section_divider(prs, '05', 'Block Diagram', divider_slide)

add_content_slide(prs, 'SHARP-LLM Pipeline \u2014 End-to-End Flow', [
    '▸  Source Code Input (C/C++ from NIST Juliet / Real codebase)',
    '          \u2193',
    '▸  Stage 1: Preprocessing & Tokenization',
    '     ◦  Header removal \u2192 Comment stripping \u2192 Guard removal \u2192 Whitespace normalization',
    '          \u2193',
    '▸  Stage 2: Fine-Tuned Lightweight LLM (CodeT5-Base / CodeBERT / GraphCodeBERT)',
    '     ◦  Encoder \u2192 Mean Pooling \u2192 Dropout \u2192 Linear \u2192 Softmax \u2192 118-dim CWE prediction',
    '          \u2193',
    '▸  Stage 3: CWE Association \u2192 Vulnerability classified into specific CWE category',
    '          \u2193',
    '▸  Stage 4: Privacy & Security Risk Categorization (R = f(c), LINDDUN mapping)',
    '          \u2193',
    '▸  Stage 5: Policy-Relevant Control Domains (C = g(R))',
    '          \u2193',
    '▸  Stage 6: Healthcare Risk Prioritization \u2192 Critical / High / Medium / Low',
], content_slide)

# ══════════════════════════════════════════════════════════════════════════
#  SECTION 6: RESULTS
# ══════════════════════════════════════════════════════════════════════════
add_section_divider(prs, '06', 'Results', divider_slide)

add_table_slide(prs, 'Model Comparison \u2014 Full Test Set (22,447 samples)', [
    'Model', 'Params', 'Accuracy', 'Precision', 'Recall', 'F1', 'MCC', 'Latency',
], [
    ['CodeT5-Small', '35.4M', '0.9980', '0.9533', '0.9534', '0.9533', '0.9979', '17.2 ms'],
    ['CodeT5-Base', '109.7M', '0.9978', '0.9585', '0.9644', '0.9605', '0.9977', '24.9 ms'],
    ['CodeBERT-Base', '124.7M', '0.9963', '0.9574', '0.9539', '0.9512', '0.9961', '19.4 ms'],
    ['GraphCodeBERT', '124.7M', '0.9975', '0.9530', '0.9530', '0.9530', '0.9974', '26.9 ms'],
    ['CodeGemma-2B (ZS)', '2.51B', '0.2500', '0.1755', '0.0832', '0.0948', '0.3022', '3,277 ms'],
    ['Gemma-4-E2B-IT (ZS)', '5.13B', '0.7800', '0.6990', '0.7167', '0.6988', '0.7799', '396,273 ms'],
], content_slide)

add_content_slide(prs, 'Key Experimental Findings', [
    '▸  CodeT5-Base is the best model: highest macro-F1 (0.9605) despite fewer params',
    '▸  Early stopping is critical: CodeBERT & GraphCodeBERT overfit at epoch 2',
    '     ◦  GraphCodeBERT F1 drops 5.78 pp (0.9530 \u2192 0.8952) at epoch 2',
    '▸  Template-aware splitting prevents inflated metrics: still >99.6% accuracy',
    '▸  Consumer-grade GPU sufficient: all encoder models train in <53 min on 8 GB VRAM',
    '▸  Fine-tuning dominates zero-shot by 28\u201375 pp accuracy and 120\u201323,000\u00d7 faster inference',
    '▸  Instruction-tuning improves zero-shot 3\u00d7 (Gemma-4 78% vs CodeGemma 25%)',
    '▸  Sub-30 ms latency \u2014 viable for CI/CD, IDE plugins, real-time code review',
], content_slide)

add_content_slide(prs, 'Per-Class Performance & Error Analysis', [
    '▸  74 of 87 test CWE categories achieve perfect F1 = 1.00 (85.1%)',
    '▸  Weighted macro-F1 = 0.9979, confirming accuracy on the vast majority of samples',
    '',
    '▸  Dominant error: misclassification toward CWE-126 (Buffer Over-read)',
    '     ◦  23 of 40 total errors are misclassified as CWE-126',
    '     ◦  CWE-127 (Buffer Under-read): 6 errors',
    '',
    '▸  Root cause: buffer access CWEs share syntactic patterns',
    '     ◦  Pointer arithmetic, array indexing, similar code structures',
    '',
    '▸  3 categories with F1 = 0.00 each have only 1 test sample',
    '     ◦  CWE-561, CWE-562, CWE-674 \u2014 insufficient data, not model failure',
], content_slide)

add_content_slide(prs, 'Training Convergence', [
    '▸  CodeT5-Small: Best at epoch 2, 53 min total \u2014 continued improving',
    '▸  CodeT5-Base: Best at epoch 2, 53 min total \u2014 continued improving',
    '▸  CodeBERT-Base: Best at epoch 1, 22 min/epoch \u2014 overfits at epoch 2',
    '     ◦  F1: 0.9512 \u2192 0.9377 (-1.35 pp)',
    '▸  GraphCodeBERT: Best at epoch 1, 23 min/epoch \u2014 overfits at epoch 2',
    '     ◦  F1: 0.9530 \u2192 0.8952 (-5.78 pp)',
    '',
    '▸  Inference Latency (speedup vs CodeGemma-2B zero-shot):',
    '     ◦  CodeT5-Small: 17.2 ms (190\u00d7 faster)',
    '     ◦  CodeT5-Base: 24.9 ms (131\u00d7 faster)',
    '     ◦  CodeBERT: 19.4 ms (169\u00d7 faster)',
    '     ◦  GraphCodeBERT: 26.9 ms (122\u00d7 faster)',
    '     ◦  Gemma-4-E2B-IT: 396,273 ms (6.6 min/sample!)',
], content_slide)

# ══════════════════════════════════════════════════════════════════════════
#  SECTION 7: REFERENCES
# ══════════════════════════════════════════════════════════════════════════
add_section_divider(prs, '07', 'References', divider_slide)

add_content_slide(prs, 'Key References (1/2)', [
    '[1]  Lu et al. (2024) \u2014 GRACE: Graph-enhanced LLM vulnerability detection, J. Sys. & SW',
    '[2]  Yang et al. (2025) \u2014 DLAP: DL-augmented LLM prompting, J. Sys. & SW',
    '[3]  Tian et al. (2025) \u2014 FUSEVUL: Multi-modal vulnerability detection, Inf. Fusion',
    '[4]  Mao et al. (2025) \u2014 LLMVulExp: Explainable vulnerability detection, IEEE TSE',
    '[5]  Liu et al. (2025) \u2014 VulACLLM: Lightweight LLM, J. Info. Sec. & Apps.',
    '[6]  Ferrag et al. (2025) \u2014 SecureFalcon: Compact LLM for vuln. classification, IEEE TSE',
    '[7]  Qiu et al. (2026) \u2014 RLV: Context-aware detection via retrieval, J. Sys. & SW',
    '[8]  Chen et al. (2026) \u2014 GPTVD: Static slicing + CoT, Auto. SW Eng.',
    '[9]  Ameh et al. (2025) \u2014 C3-VULMAP: Privacy-aware healthcare vuln. dataset, Electronics',
    '[10] NSA/NIST (2017) \u2014 Juliet Test Suite for C/C++ v1.3',
], content_slide)

add_content_slide(prs, 'Key References (2/2)', [
    '[11] Wang et al. (2021) \u2014 CodeT5: Identifier-aware Unified Pre-trained Encoder-Decoder, EMNLP',
    '[12] Feng et al. (2020) \u2014 CodeBERT: Pre-trained Model for Programming Languages, EMNLP',
    '[13] Guo et al. (2020) \u2014 GraphCodeBERT: Pre-training Code Representations with Data Flow',
    '[14] Dettmers et al. (2023) \u2014 QLoRA: Efficient Finetuning of Quantized LLMs, NeurIPS',
    '[15] CodeGemma Team (2024) \u2014 CodeGemma: Open Code Models',
    '[16] Li et al. (2018) \u2014 VulDeePecker: Deep learning-based vulnerability detection, NDSS',
    '[17] Fan et al. (2020) \u2014 Big-Vul: Large-scale vulnerability dataset, MSR',
    '[18] Zhou et al. (2019) \u2014 Devign: Effective vulnerability identification, NeurIPS',
    '[19] Dolcetti & Iotti (2025) \u2014 Dual perspective review on LLMs and code verification, Frontiers',
    '[20] Shaon & Akter (2025) \u2014 Modern approaches to software vulnerability detection, Electronics',
], content_slide)

# ── THANK YOU SLIDE ──────────────────────────────────────────────────────
# Add a Thank You slide using the divider template
thank_you = _clone_slide(prs, divider_slide)
# Update to show "Thank You" instead of a number
for shape in thank_you.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        txt = ''.join(r.text for r in para.runs).strip()
        if txt in ('01', '02', '03', '04', '05', '06', '07', '08'):
            for run in para.runs:
                run.text = ''
            # Hide or blank the number
        for run in para.runs:
            if run.font.bold and run.font.size and run.font.size >= Pt(28) and len(txt) > 3:
                for r in para.runs:
                    r.text = ''
                para.runs[0].text = 'Thank You'
                break

# ── Save ─────────────────────────────────────────────────────────────────
prs.save(DST)
print(f'Saved updated presentation to:\n{DST}')
print(f'Total slides: {len(prs.slides)}')
