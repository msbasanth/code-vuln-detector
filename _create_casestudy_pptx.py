"""Generate the DL Case Study presentation for SHARP-LLM.

Uses the base presentation as a template to preserve theme, logos, and layouts.
Follows the guideline structure:
1. Title Slide
2. Introduction
3. Problem Statement
4. Objectives
5. Methodology
6. Case Study Details
7. Analysis
8. Results / Outcomes
9. Discussion
10. Conclusion
11. References
"""

import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

# ── Load base presentation and extract theme / layouts ──────────────────────
BASE_PATH = "paper/SHARP-LLM_A_Framework_for_Vulnerabilitty_Detection_AM.SC.R4CSE25007-PT.pptx"
OUTPUT_PATH = "paper/SHARP-LLM_DL_CaseStudy_Presentation.pptx"

base_prs = Presentation(BASE_PATH)

# Identify layouts
layout_map = {sl.name: sl for sl in base_prs.slide_layouts}
title_layout = layout_map.get("Title Slide")
content_layout = layout_map.get("4_Custom Layout 1 1")  # Main content layout
title_card_layout = layout_map.get("4_Custom Layout 1")  # Title card layout (slide 2 style)
section_layout = content_layout  # Section headers use same layout

# ── Helper functions ────────────────────────────────────────────────────────

def clone_slide(prs, template_slide_index):
    """Clone a slide from the presentation (preserving all shapes, images, backgrounds)."""
    template = prs.slides[template_slide_index]
    slide_layout = template.slide_layout
    slide = prs.slides.add_slide(slide_layout)

    # Copy background
    bg_elem = template.background._element
    slide_bg = slide.background._element
    for child in list(slide_bg):
        slide_bg.remove(child)
    for child in bg_elem:
        slide_bg.append(copy.deepcopy(child))

    return slide


def clear_slide_shapes(slide):
    """Remove all shapes from a slide (keep background)."""
    sp_tree = slide.shapes._spTree
    for sp in list(sp_tree):
        tag = etree.QName(sp.tag).localname if isinstance(sp.tag, str) else ""
        if tag in ("sp", "pic", "graphicFrame", "grpSp"):
            sp_tree.remove(sp)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if color:
        p.font.color.rgb = color
    return txBox


def add_bullet_slide(slide, title_text, bullets, left=Inches(0.8), top=Inches(1.6),
                     width=Inches(8.5), height=Inches(4.5), font_size=16):
    """Add title + bullet points to a content slide."""
    # Title
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(9), Inches(0.8),
                title_text, font_size=28, bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))

    # Bullets
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle sub-bullets (indented with spaces or ◦)
        if bullet.startswith("  ") or bullet.startswith("◦"):
            clean = bullet.lstrip(" ◦").strip()
            p.text = f"◦  {clean}"
            p.level = 1
            p.font.size = Pt(font_size - 2)
        else:
            clean = bullet.lstrip("▸ ").strip()
            p.text = f"▸  {clean}"
            p.level = 0
            p.font.size = Pt(font_size)

        p.font.name = "Calibri"
        p.space_after = Pt(4)
        p.space_before = Pt(2)

    return txBox


def add_section_slide(prs, section_number, section_title):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(content_layout)
    clear_slide_shapes(slide)

    # Section number
    add_textbox(slide, Inches(3.5), Inches(2.0), Inches(3), Inches(1.2),
                f"{section_number:02d}", font_size=72, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C), alignment=PP_ALIGN.CENTER)

    # Section title
    add_textbox(slide, Inches(1.5), Inches(3.2), Inches(7), Inches(1.0),
                section_title, font_size=32, bold=False,
                color=RGBColor(0x4A, 0x4A, 0x4A), alignment=PP_ALIGN.CENTER)

    return slide


def add_table_slide(slide, title_text, headers, rows, col_widths=None,
                    left=Inches(0.5), top=Inches(1.6)):
    """Add a slide with a title and table."""
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(9), Inches(0.8),
                title_text, font_size=28, bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))

    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)
    width = Inches(9)
    height = Inches(0.35 * n_rows)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.name = "Calibri"
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = "Calibri"
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)

    return table_shape


# ── Build the new presentation ──────────────────────────────────────────────
# We'll delete all existing slides and rebuild from scratch using the same layouts

# First, remove all slides except the first one (keep slide master/theme)
# Actually, let's create a fresh presentation from the template
prs = Presentation(BASE_PATH)

# Remove all existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId is None:
        # Try the 'r:id' attribute
        rId = prs.slides._sldIdLst[0].attrib.get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# Re-map layouts from the new prs
layout_map = {sl.name: sl for sl in prs.slide_layouts}
title_layout = layout_map.get("Title Slide")
content_layout = layout_map.get("4_Custom Layout 1 1")
title_card_layout = layout_map.get("4_Custom Layout 1")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide (University branding)
# ════════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(title_layout)
clear_slide_shapes(slide1)
add_textbox(slide1, Inches(1.0), Inches(1.5), Inches(8), Inches(0.6),
            "Deep Learning Case Study", font_size=18,
            color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)
add_textbox(slide1, Inches(0.5), Inches(2.1), Inches(9), Inches(0.6),
            "Amrita School of Computing", font_size=14,
            color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)
add_textbox(slide1, Inches(0.5), Inches(2.5), Inches(9), Inches(0.5),
            "Amritapuri Campus", font_size=12,
            color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Title Card
# ════════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(title_card_layout)
clear_slide_shapes(slide2)
add_textbox(slide2, Inches(0.8), Inches(1.0), Inches(8.5), Inches(1.2),
            "SHARP-LLM", font_size=44, bold=True,
            color=RGBColor(0x1B, 0x3A, 0x5C), alignment=PP_ALIGN.LEFT)
add_textbox(slide2, Inches(0.8), Inches(2.0), Inches(8.5), Inches(1.0),
            "A Secure and Resource-Efficient Framework for Source Code\n"
            "Vulnerability Analysis and Healthcare-Specific Risk Prioritization",
            font_size=20, color=RGBColor(0x4A, 0x4A, 0x4A))
add_textbox(slide2, Inches(0.8), Inches(3.5), Inches(8.5), Inches(0.5),
            "Basanth M S (AM.SC.R4CSE25007-PT),  Gopakumar G",
            font_size=16, bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
add_textbox(slide2, Inches(0.8), Inches(4.1), Inches(8.5), Inches(0.8),
            "Dept. of Computer Science and Engineering, School of Computing\n"
            "Amrita Vishwa Vidyapeetham, Amritapuri, India",
            font_size=14, color=RGBColor(0x66, 0x66, 0x66))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Section — Introduction
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, 1, "Introduction")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Introduction Content
# ════════════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide4)
add_bullet_slide(slide4, "Background and Context", [
    "Software vulnerabilities cause security breaches, service disruption, and data exposure",
    "Healthcare software is especially critical — vulnerabilities can compromise:",
    "  Electronic Health Records (EHR) and patient data confidentiality",
    "  Clinical decision support and diagnostic workflows",
    "  Telemedicine platforms and connected medical devices",
    "  Regulatory compliance (HIPAA, GDPR)",
    "Traditional SAST tools have high false-positive rates and limited generalization",
    "LLM-based approaches show promise but are often computationally expensive",
    "No existing framework combines lightweight detection with healthcare-oriented risk prioritization",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Section — Problem Statement
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, 2, "Problem Statement")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Problem Statement Content
# ════════════════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide6)
add_bullet_slide(slide6, "Problem Statement", [
    "Source code vulnerability detection in C/C++ programs using fine-tuned lightweight LLMs",
    "Given a code sample, determine whether it contains vulnerability-relevant patterns",
    "  Classify into one of 118 CWE categories",
    "Additionally, prioritize findings in healthcare contexts where defects may impact:",
    "  Patient data confidentiality and integrity",
    "  Service continuity and system trust",
    "  Regulatory compliance",
    "Three key challenges:",
    "  Vulnerable vs. safe code may differ in only minor semantic details",
    "  118 CWE categories with uneven distributions (1,500 vs. 5 samples)",
    "  High-performing models must be practical on consumer hardware",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Section — Objectives
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, 3, "Objectives")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Objectives Content
# ════════════════════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide8)
add_bullet_slide(slide8, "What the Study Aims to Achieve", [
    "Evaluate lightweight LLMs for 118-class CWE vulnerability classification on C/C++ code",
    "Compare three model paradigms on the same benchmark:",
    "  Encoder fine-tuning (CodeT5, CodeBERT, GraphCodeBERT)",
    "  QLoRA fine-tuning (CodeGemma-2B, 4-bit quantized)",
    "  Zero-shot generative inference (CodeGemma-2B, Gemma-4-E2B-IT)",
    "Integrate a BiLSTM with Self-Attention as a traditional DL baseline",
    "  ~6.5M parameters — benchmarks transformer models against classical DL",
    "Introduce a healthcare-oriented risk prioritization layer:",
    "  CWE → LINDDUN privacy threats → Control domains → Risk score",
    "Demonstrate practical deployment on consumer-grade GPU (8 GB VRAM, <30 ms latency)",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Section — Methodology
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, 4, "Methodology")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10: Dataset & Preprocessing
# ════════════════════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide10)
add_bullet_slide(slide10, "Dataset & Preprocessing", [
    "NIST Juliet Test Suite for C/C++ v1.3",
    "  105,183 source code samples across 118 CWE categories",
    "  Training: 82,736 samples (1,323 templates)",
    "  Test: 22,447 samples (366 templates, 87 CWEs present)",
    "  Zero template overlap (GroupShuffleSplit 80/20, seed=42)",
    "4-Stage Preprocessing Pipeline:",
    "  Stage 1: Header removal (template-generated comment blocks)",
    "  Stage 2: Comment stripping (state machine, preserves string literals)",
    "  Stage 3: Guard removal (#ifndef OMITBAD / OMITGOOD)",
    "  Stage 4: Whitespace normalization (compact representation)",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11: Models Evaluated (Table)
# ════════════════════════════════════════════════════════════════════════════
slide11 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide11)
add_table_slide(slide11, "Models Evaluated", 
    ["Model", "Paradigm", "Params", "Key Details"],
    [
        ["CodeT5-Small", "Encoder FT", "35.4M", "Full fine-tuning, FP16"],
        ["CodeT5-Base", "Encoder FT", "109.7M", "Full fine-tuning, FP16"],
        ["CodeBERT-Base", "Encoder FT", "124.7M", "Full fine-tuning, FP16"],
        ["GraphCodeBERT", "Encoder FT", "124.7M", "Full fine-tuning, FP16"],
        ["BiLSTM-Attention", "DL Baseline", "6.5M", "Embedding→BiLSTM→Self-Attn→Linear"],
        ["CodeGemma-2B", "QLoRA FT", "2.51B", "4-bit NF4, LoRA r=16, α=32"],
        ["CodeGemma-2B", "Zero-Shot", "2.51B", "Beam search B=10"],
        ["Gemma-4-E2B-IT", "Zero-Shot", "5.13B", "Instruction-tuned, B=2"],
    ])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12: Training Configuration
# ════════════════════════════════════════════════════════════════════════════
slide12 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide12)
add_bullet_slide(slide12, "Training Configuration", [
    "Encoder Models (CodeT5, CodeBERT, GraphCodeBERT):",
    "  LR: 5×10⁻⁵, AdamW, batch 8, max 256 tokens, 2 epochs, FP16, patience=1",
    "BiLSTM-Attention (DL Baseline):",
    "  LR: 1×10⁻³, AdamW, batch 32, max 256 tokens, 10 epochs, FP16, patience=3",
    "  Embedding(32K, 128) → BiLSTM(128→256, 2 layers) → Self-Attention → Linear(512→118)",
    "  Reuses CodeT5's BPE tokenizer for consistent tokenization",
    "QLoRA (CodeGemma-2B):",
    "  LR: 2×10⁻⁴, PagedAdamW 8-bit, LoRA targets: q/k/v/o/gate/up/down",
    "Hardware: NVIDIA RTX 2000 Ada (8 GB VRAM), Intel i7-13850HX",
    "Stack: Python 3.13, PyTorch ≥2.0, HuggingFace Transformers ≥5.5",
], font_size=15)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13: Healthcare Risk Prioritization
# ════════════════════════════════════════════════════════════════════════════
slide13 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide13)
add_bullet_slide(slide13, "Healthcare Risk Prioritization Layer", [
    "Step 1 — Privacy & Security Risk Categorization:",
    "  R = f(c) — map detected CWE to LINDDUN threat types",
    "  Linkability, Identifiability, Non-repudiation, Detectability,",
    "  Information Disclosure, Unawareness, Non-compliance",
    "Step 2 — Policy-Relevant Control Domains:",
    "  C = g(R) — confidentiality, access control, audit & accountability,",
    "  secure data handling, service availability",
    "Step 3 — Weighted Risk Prioritization:",
    "  P = w₁·s + w₂·E + w₃·D + w₄·S + w₅·C'",
    "  s = model confidence, E = severity, D = data sensitivity",
    "  S = service/safety impact, C' = control domain importance",
    "Output: Ordinal priority level (Critical / High / Medium / Low)",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14: Section — Case Study Details
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, 5, "Case Study Details")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15: Case Study Details
# ════════════════════════════════════════════════════════════════════════════
slide15 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide15)
add_bullet_slide(slide15, "Experimental Design — 7 Experiments", [
    "Experiment A: Juliet-118 Baseline — 118 CWEs, 82K train / 22K test, 4 encoder models",
    "Experiment B: Juliet-19 Subset — 19 overlapping CWEs with Big-Vul",
    "Experiment C: Juliet + Big-Vul Combined — 19 CWEs, real-world code added",
    "Experiment D: Big-Vul Only — 55 CWEs from real-world vulnerability dataset",
    "Experiment E: Juliet + Big-Vul Union — 187 CWEs, full coverage",
    "Experiment F: Union Extended — 6 epochs, patience=3 (vs 2 ep in Exp E)",
    "Experiment G: BiLSTM-Attention DL Baseline — 118 CWEs, 10 epochs",
    "  Traditional DL model to benchmark against transformer-based approaches",
    "  Tests whether simpler architecture can match LLM performance on structured data",
], font_size=15)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16: Three-fold Novelty
# ════════════════════════════════════════════════════════════════════════════
slide16 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide16)
add_bullet_slide(slide16, "Three-Fold Novelty", [
    "Novelty 1: Template-Aware Data Splitting",
    "  Eliminates structural data leakage inherent in the Juliet Test Suite",
    "  Splits at template level using GroupShuffleSplit — zero template overlap",
    "  First rigorous anti-leakage split for Juliet-based evaluations",
    "Novelty 2: Multi-Paradigm Model Evaluation",
    "  First study comparing 4 paradigms on 118-class CWE classification:",
    "  Encoder FT | QLoRA FT | Zero-shot generative | Traditional DL (BiLSTM)",
    "Novelty 3: Healthcare-Oriented Risk Prioritization Layer",
    "  Three-stage post-detection pipeline:",
    "  CWE → LINDDUN threat types → Policy control domains → Weighted risk score",
    "  P = λ₁s + λ₂E + λ₃D + λ₄S + λ₅C' → Priority Level",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17: Section — Analysis
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, 6, "Analysis")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18: Model Comparison Table
# ════════════════════════════════════════════════════════════════════════════
slide18 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide18)
add_table_slide(slide18, "Model Comparison — Experiment A (Juliet-118, Full Test Set)",
    ["Model", "Params", "Acc", "Prec", "Recall", "F1", "Latency"],
    [
        ["CodeT5-Small", "35.4M", "0.9972", "0.9349", "0.9410", "0.9370", "17.2 ms"],
        ["CodeT5-Base", "109.7M", "0.9977", "0.9309", "0.9418", "0.9345", "24.9 ms"],
        ["CodeBERT-Base", "124.7M", "0.9974", "0.9477", "0.9528", "0.9493", "19.4 ms"],
        ["GraphCodeBERT", "124.7M", "0.9975", "0.9530", "0.9530", "0.9530", "26.9 ms"],
        ["BiLSTM-Attention", "6.5M", "0.9982", "0.9997", "0.9990", "0.9993", "—"],
        ["CodeGemma (ZS)", "2.51B", "0.2500", "0.1755", "0.0832", "0.0948", "3,277 ms"],
        ["Gemma-4 IT (ZS)", "5.13B", "0.7800", "0.6990", "0.7167", "0.6988", "396s"],
    ])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 19: Cross-Experiment Results
# ════════════════════════════════════════════════════════════════════════════
slide19 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide19)
add_table_slide(slide19, "Cross-Experiment Macro-F1 Comparison",
    ["Experiment", "CodeT5-S", "CodeT5-B", "CodeBERT", "GraphCB", "BiLSTM"],
    [
        ["A: Juliet-118", "0.9370", "0.9345", "0.9493", "0.9530", "0.9993"],
        ["B: Juliet-19", "0.9470", "0.9367", "0.9987", "0.9472", "—"],
        ["C: Combined-19", "0.8239", "0.8171", "0.8011", "0.8310", "—"],
        ["D: Big-Vul only", "0.0462", "0.0656", "0.0074", "0.0730", "—"],
        ["E: Union-187", "0.5974", "0.6045", "0.5901", "0.6219", "—"],
        ["F: Union-6ep", "0.6376", "0.6312", "0.6142", "0.6292", "—"],
        ["G: BiLSTM-118", "—", "—", "—", "—", "0.9993"],
    ],
    top=Inches(1.5))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 20: Section — Results / Outcomes
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, 7, "Results / Outcomes")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 21: Key Results
# ════════════════════════════════════════════════════════════════════════════
slide21 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide21)
add_bullet_slide(slide21, "Key Experimental Findings", [
    "BiLSTM-Attention achieves highest Macro-F1 (0.9993) on Juliet-118 with only 6.5M params",
    "  19× fewer parameters than GraphCodeBERT (124.7M), yet outperforms all transformers",
    "GraphCodeBERT is best transformer: F1=0.9530 (Exp A), but overfits at epoch 2",
    "  F1 drops 5.78 pp (0.9530 → 0.8952) — early stopping is critical",
    "Fine-tuning dominates zero-shot by 28–75 pp accuracy and 120–23,000× faster inference",
    "  Instruction-tuning improves zero-shot 3× (Gemma-4 78% vs CodeGemma 25%)",
    "Real-world data (Big-Vul) severely degrades all models: F1 drops to 0.01–0.07",
    "  Gap between synthetic (Juliet) and real-world vulnerability patterns",
    "Consumer-grade GPU sufficient: all models train in <53 min on 8 GB VRAM",
    "Sub-30 ms latency — viable for CI/CD, IDE plugins, real-time code review",
], font_size=15)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 22: Per-Class Performance
# ════════════════════════════════════════════════════════════════════════════
slide22 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide22)
add_bullet_slide(slide22, "Per-Class Performance & Error Analysis", [
    "74 of 87 test CWE categories achieve perfect F1 = 1.00 (85.1%)",
    "Weighted macro-F1 = 0.9982, confirming accuracy on the vast majority of samples",
    "BiLSTM confusion pattern: misclassifications toward CWE-176 (Unicode Encoding)",
    "  8 errors from CWE-114, 6 from CWE-194, 4 from CWE-121/123/126/366",
    "Transformer confusion pattern: misclassifications toward CWE-126 (Buffer Over-read)",
    "  Buffer access CWEs share syntactic patterns (pointer arithmetic, array indexing)",
    "3 categories with F1 = 0.00 each have only 1 test sample",
    "  CWE-561, CWE-562, CWE-674 — insufficient data, not model failure",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 23: Section — Discussion
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, 8, "Discussion")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 24: Discussion Content
# ════════════════════════════════════════════════════════════════════════════
slide24 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide24)
add_bullet_slide(slide24, "Discussion — Insights & Implications", [
    "BiLSTM vs Transformers on structured synthetic data:",
    "  BiLSTM excels when patterns are learnable from token sequences alone",
    "  Transformers may be over-parameterized for highly structured Juliet data",
    "  Real-world generalization (Exp D) remains the key challenge for all models",
    "Healthcare risk prioritization adds actionable context:",
    "  Same CWE has different priority in EHR system vs game engine",
    "  LINDDUN mapping provides privacy-aware vulnerability triage",
    "Practical deployment considerations:",
    "  6.5M-param BiLSTM is ideal for edge/embedded healthcare devices",
    "  Transformer models offer better generalization potential on diverse codebases",
    "  Framework supports CI/CD integration with sub-30 ms inference",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 25: Related Work Comparison
# ════════════════════════════════════════════════════════════════════════════
slide25 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide25)
add_table_slide(slide25, "Comparison with Related Work",
    ["Framework", "Focus", "Key Limitation", "SHARP-LLM Advantage"],
    [
        ["GRACE", "Graph+LLM+ICL", "Prompt sensitivity", "No prompting needed"],
        ["DLAP", "DL-guided prompting", "DL model dependency", "End-to-end pipeline"],
        ["FUSEVUL", "Multi-modal", "Basic prompting", "Healthcare prioritization"],
        ["VulACLLM", "AST+CFG+LoRA", "Compression trade-off", "Lightweight by design"],
        ["SecureFalcon", "Compact LLM", "Limited patterns", "118 CWE + risk layer"],
        ["GPTVD", "Static slicing+CoT", "Scalability", "Consumer GPU feasible"],
    ],
    top=Inches(1.5))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 26: Section — Conclusion
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, 9, "Conclusion")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 27: Conclusion Content
# ════════════════════════════════════════════════════════════════════════════
slide27 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide27)
add_bullet_slide(slide27, "Summary & Key Takeaways", [
    "SHARP-LLM: lightweight, secure framework for vulnerability detection + healthcare risk prioritization",
    "Key results on Juliet-118 (22,447 test samples):",
    "  BiLSTM-Attention: Macro-F1 = 0.9993 (6.5M params) — best overall",
    "  GraphCodeBERT: Macro-F1 = 0.9530 (124.7M params) — best transformer",
    "  All encoder models: >99.7% accuracy, <30 ms latency",
    "Three-fold novelty: template-aware splitting, multi-paradigm comparison, healthcare risk layer",
    "Limitations and future work:",
    "  Real-world generalization gap (Big-Vul F1 < 0.07) — needs data augmentation",
    "  BiLSTM high F1 on Juliet may reflect structural regularity of synthetic data",
    "  Extend to more languages (Java, Python) and real-world healthcare codebases",
    "  Integrate into VS Code extension or CI/CD pipeline for practical deployment",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 28: Section — References
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, 10, "References")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 29: References 1/2
# ════════════════════════════════════════════════════════════════════════════
slide29 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide29)
add_bullet_slide(slide29, "Key References (1/2)", [
    "[1]  Lu et al. (2024) — GRACE: Graph-enhanced LLM vulnerability detection, J. Sys. & SW",
    "[2]  Yang et al. (2025) — DLAP: DL-augmented LLM prompting, J. Sys. & SW",
    "[3]  Tian et al. (2025) — FUSEVUL: Multi-modal vulnerability detection, Inf. Fusion",
    "[4]  Mao et al. (2025) — LLMVulExp: Explainable vulnerability detection, IEEE TSE",
    "[5]  Liu et al. (2025) — VulACLLM: Lightweight LLM, J. Info. Sec. & Apps.",
    "[6]  Ferrag et al. (2025) — SecureFalcon: Compact LLM for vuln. classification, IEEE TSE",
    "[7]  Qiu et al. (2026) — RLV: Context-aware detection via retrieval, J. Sys. & SW",
    "[8]  Chen et al. (2026) — GPTVD: Static slicing + CoT, Auto. SW Eng.",
    "[9]  Ameh et al. (2025) — C3-VULMAP: Privacy-aware healthcare vuln. dataset, Electronics",
    "[10] NSA/NIST (2017) — Juliet Test Suite for C/C++ v1.3",
], font_size=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 30: References 2/2
# ════════════════════════════════════════════════════════════════════════════
slide30 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide30)
add_bullet_slide(slide30, "Key References (2/2)", [
    "[11] Wang et al. (2021) — CodeT5: Identifier-aware Pre-trained Encoder-Decoder, EMNLP",
    "[12] Feng et al. (2020) — CodeBERT: Pre-trained Model for Programming Languages, EMNLP",
    "[13] Guo et al. (2020) — GraphCodeBERT: Pre-training Code Representations with Data Flow",
    "[14] Dettmers et al. (2023) — QLoRA: Efficient Finetuning of Quantized LLMs, NeurIPS",
    "[15] CodeGemma Team (2024) — CodeGemma: Open Code Models",
    "[16] Li et al. (2018) — VulDeePecker: Deep learning-based vulnerability detection, NDSS",
    "[17] Fan et al. (2020) — Big-Vul: Large-scale vulnerability dataset, MSR",
    "[18] Zhou et al. (2019) — Devign: Effective vulnerability identification, NeurIPS",
    "[19] Dolcetti & Iotti (2025) — Dual perspective review on LLMs and code verification, Frontiers",
    "[20] Shaon & Akter (2025) — Modern approaches to software vulnerability detection, Electronics",
], font_size=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 31: Thank You
# ════════════════════════════════════════════════════════════════════════════
slide31 = prs.slides.add_slide(content_layout)
clear_slide_shapes(slide31)
add_textbox(slide31, Inches(1), Inches(2.5), Inches(8), Inches(1.5),
            "Thank You", font_size=48, bold=True,
            color=RGBColor(0x1B, 0x3A, 0x5C), alignment=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────────────────
prs.save(OUTPUT_PATH)
print(f"Presentation saved to: {OUTPUT_PATH}")
print(f"Total slides: {len(prs.slides)}")
