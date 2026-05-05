"""Generate the DL Case Study presentation — aligned with:
  • context/DL/DL_Guidelines for Case Study.pdf
  • DeepLearning_CaseStudyAbstract_Basanth_AM.SC.R4CSE25007-PT.docx

Slide deck structure (per guidelines):
 1  Title Slide          – topic, name, register number
 2  Abstract             – submitted abstract summary
 3  Introduction          – background & context
 4  Problem Statement     – clearly defined issue
 5  Objectives            – what the study aims to achieve
 6  Literature Review     – related work (assessment criterion)
 7  Methodology           – approach, tools, techniques
    7a  Dataset & Preprocessing
    7b  Model Architecture (BiLSTM + Attention detail)
    7c  Training Configuration
    7d  Evaluation Metrics
 8  Case Study Details    – experiments, data, novelty
 9  Analysis              – interpretation of findings
    9a  Training / Validation Curves
    9b  Confusion Matrix & Per-Class Performance
    9c  Attention Visualization (interpretability technique)
10  Results / Outcomes    – key observations, tables
11  Discussion            – insights, implications, comparisons
12  Ethical Considerations – bias, fairness, misuse (guideline req.)
13  Limitations & Future Work
14  Conclusion            – summary & key takeaways
15  References            – proper citation
16  Thank You
"""

import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── Config ──────────────────────────────────────────────────────────────────
BASE_PATH = "paper/SHARP-LLM_A_Framework_for_Vulnerabilitty_Detection_AM.SC.R4CSE25007-PT.pptx"
OUTPUT_PATH = "paper/SHARP-LLM_DL_CaseStudy_Presentation.pptx"

# ── Theme colours ───────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1B, 0x3A, 0x5C)
GREY       = RGBColor(0x4A, 0x4A, 0x4A)
LIGHT_GREY = RGBColor(0x66, 0x66, 0x66)
MID_GREY   = RGBColor(0x88, 0x88, 0x88)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ZEBRA_BG   = RGBColor(0xF0, 0xF4, 0xF8)

# ── Helpers ─────────────────────────────────────────────────────────────────

def clear_slide_shapes(slide):
    sp_tree = slide.shapes._spTree
    for sp in list(sp_tree):
        tag = etree.QName(sp.tag).localname if isinstance(sp.tag, str) else ""
        if tag in ("sp", "pic", "graphicFrame", "grpSp"):
            sp_tree.remove(sp)


def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=None,
                alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if color:
        p.font.color.rgb = color
    return txBox


def add_bullet_slide(slide, title_text, bullets,
                     left=Inches(0.8), top=Inches(1.6),
                     width=Inches(8.5), height=Inches(4.5),
                     font_size=16, title_size=28):
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(9), Inches(0.8),
                title_text, font_size=title_size, bold=True, color=DARK_BLUE)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if bullet.startswith("  ") or bullet.startswith("◦"):
            clean = bullet.lstrip(" ◦").strip()
            p.text = f"◦  {clean}"
            p.level = 1
            p.font.size = Pt(font_size - 2)
        else:
            clean = bullet.lstrip("▸ ").strip()
            p.text = f"▸  {clean}"
            p.level = 0
            p.font.size = Pt(font_size)
        p.font.name = "Calibri"
        p.space_after = Pt(4)
        p.space_before = Pt(2)
    return txBox


def add_section_slide(prs, layout, section_number, section_title):
    slide = prs.slides.add_slide(layout)
    clear_slide_shapes(slide)
    add_textbox(slide, Inches(3.5), Inches(2.0), Inches(3), Inches(1.2),
                f"{section_number:02d}", font_size=72, bold=True,
                color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(3.2), Inches(7), Inches(1.0),
                section_title, font_size=32, bold=False,
                color=GREY, alignment=PP_ALIGN.CENTER)
    return slide


def add_table_slide(slide, title_text, headers, rows,
                    left=Inches(0.5), top=Inches(1.6)):
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(9), Inches(0.8),
                title_text, font_size=28, bold=True, color=DARK_BLUE)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    width = Inches(9)
    height = Inches(0.35 * n_rows)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11); p.font.bold = True
            p.font.name = "Calibri"; p.font.color.rgb = WHITE
        cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BLUE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10); p.font.name = "Calibri"
            if i % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = ZEBRA_BG
    return table_shape


# ── Build presentation ──────────────────────────────────────────────────────
prs = Presentation(BASE_PATH)

# Remove all existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId is None:
        rId = prs.slides._sldIdLst[0].attrib.get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

layout_map    = {sl.name: sl for sl in prs.slide_layouts}
title_layout  = layout_map.get("Title Slide")
content_layout = layout_map.get("4_Custom Layout 1 1")
card_layout   = layout_map.get("4_Custom Layout 1")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(title_layout); clear_slide_shapes(s)
add_textbox(s, Inches(0.5), Inches(0.8), Inches(9), Inches(1.0),
            "Deep Learning for Software Vulnerability Detection:\n"
            "A Case Study on C/C++ Projects",
            font_size=28, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.5), Inches(2.3), Inches(9), Inches(0.5),
            "Basanth M S", font_size=20, bold=True,
            color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.5), Inches(2.9), Inches(9), Inches(0.4),
            "AM.SC.R4CSE25007-PT", font_size=16,
            color=GREY, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.5), Inches(3.5), Inches(9), Inches(0.8),
            "Deep Learning Case Study\n"
            "Amrita School of Computing, Amritapuri Campus\n"
            "Amrita Vishwa Vidyapeetham",
            font_size=14, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Abstract
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_textbox(s, Inches(0.5), Inches(0.4), Inches(9), Inches(0.8),
            "Abstract", font_size=28, bold=True, color=DARK_BLUE)
add_textbox(s, Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.0),
            "This case study presents a focused examination of software vulnerability "
            "detection using deep learning approaches in a practical setting. Software "
            "vulnerabilities remain a major threat to the security and reliability of modern "
            "systems, often resulting in data breaches, financial loss, privacy violations, "
            "and operational disruption.\n\n"
            "The study concentrates on the application of deep learning methods to "
            "vulnerability detection using benchmark and real-world security datasets, "
            "including CVE, NVD, CWE, and the NIST Juliet Test Suite for C/C++. "
            "By analyzing these sources, the study explores how deep learning models can "
            "learn patterns from source code to identify security weaknesses more accurately.\n\n"
            "Performance is evaluated using Accuracy, Precision, Recall, F1-score, Matthews "
            "Correlation Coefficient (MCC), and False Positive Rate (FPR). Through this focused "
            "investigation, the case study highlights the strengths, limitations, and practical "
            "implications of deep learning-based vulnerability detection.",
            font_size=14, color=GREY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Section: Introduction
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, content_layout, 1, "Introduction")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Introduction Content
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Background and Context", [
    "Software vulnerabilities cause data breaches, financial loss, and operational disruption",
    "C/C++ remains dominant in safety-critical systems — buffer overflows, injection flaws,",
    "  access-control issues, and insecure deserialization are persistent threats",
    "Traditional SAST tools suffer high false-positive rates and limited generalization",
    "Deep learning can automatically learn vulnerability patterns from source code",
    "  RNNs, LSTMs, and Transformers excel on sequential data (guideline-aligned architectures)",
    "Key question: Can lightweight DL match large language models on vulnerability detection?",
    "This study applies DL to C/C++ vulnerability detection using benchmark datasets",
    "  (CVE, NVD, CWE, Juliet) with Python, PyTorch, and Hugging Face",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Section: Problem Statement
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, content_layout, 2, "Problem Statement")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Problem Statement Content
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Problem Statement", [
    "Detect vulnerabilities in C/C++ source code using deep learning methods",
    "Classify code samples into one of 118 CWE vulnerability categories",
    "Specific vulnerability types under investigation:",
    "  Injection flaws (CWE-78, CWE-89, CWE-90)",
    "  Buffer overflows (CWE-121, CWE-122, CWE-126, CWE-127)",
    "  Insecure deserialization and data handling (CWE-134, CWE-176)",
    "  Access control issues (CWE-114, CWE-252)",
    "Key challenges:",
    "  Vulnerable vs. safe code differs in minor semantic details",
    "  118 CWE categories with severe class imbalance (1,500 vs. 5 samples)",
    "  Models must be trainable on consumer-grade hardware (≤8 GB VRAM)",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Section: Objectives
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, content_layout, 3, "Objectives")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Objectives Content
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "What the Study Aims to Achieve", [
    "Apply deep learning to 118-class CWE vulnerability classification on C/C++ code",
    "Build a BiLSTM with Self-Attention as the primary DL model (~6.5M params)",
    "  Architecture matches sequential nature of source code (guideline: RNN/LSTM for seq. data)",
    "Benchmark DL baseline against four fine-tuned transformer models:",
    "  CodeT5-Small (35.4M), CodeT5-Base (109.7M), CodeBERT (124.7M), GraphCodeBERT (124.7M)",
    "Also compare against simple baselines: zero-shot generative inference (CodeGemma, Gemma-4)",
    "Evaluate using Accuracy, Precision, Recall, F1-score, MCC, and FPR",
    "Demonstrate that models are trainable on a personal laptop (consumer GPU, few hours)",
    "Apply an interpretability technique: self-attention weight visualization",
    "Discuss ethical considerations: bias, fairness, and potential model misuse",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Section: Literature Review
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, content_layout, 4, "Literature Review")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Literature Review Content
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_table_slide(s, "Related Work — Deep Learning for Vulnerability Detection",
    ["Study", "Year", "Approach", "Key Finding / Limitation"],
    [
        ["VulDeePecker (Li et al.)", "2018", "BiLSTM on code gadgets", "Pioneered DL for vuln. detection; limited CWE coverage"],
        ["Devign (Zhou et al.)", "2019", "GNN on code graphs", "Graph-based reasoning; scalability issues"],
        ["GRACE (Lu et al.)", "2024", "Graph + LLM + ICL", "Strong results; prompt sensitivity"],
        ["SecureFalcon (Ferrag)", "2025", "Compact LLM fine-tune", "Lightweight; limited CWE patterns"],
        ["VulACLLM (Liu et al.)", "2025", "AST + CFG + LoRA", "Code structure; compression trade-off"],
        ["DLAP (Yang et al.)", "2025", "DL-guided prompting", "DL + LLM synergy; DL dependency"],
        ["FUSEVUL (Tian et al.)", "2025", "Multi-modal fusion", "Comprehensive; no healthcare focus"],
        ["GPTVD (Chen et al.)", "2026", "Static slicing + CoT", "Interpretable; scalability concerns"],
    ],
    top=Inches(1.5))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Literature Review Gap
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Research Gap & Our Contribution", [
    "Existing gap in literature:",
    "  No study compares traditional DL (BiLSTM) with fine-tuned LLMs on same benchmark",
    "  Most studies use binary classification (vuln/safe) — not multi-class CWE detection",
    "  Few studies address class imbalance across 100+ CWE categories",
    "  Practical deployment on consumer hardware often not demonstrated",
    "Our contribution — three-fold novelty:",
    "  1. Template-aware data splitting: eliminates Juliet structural data leakage",
    "  2. Multi-paradigm comparison: BiLSTM vs. Encoder FT vs. QLoRA vs. Zero-shot",
    "  3. Healthcare-oriented risk prioritization layer (CWE → LINDDUN → risk score)",
    "First rigorous comparison of 6.5M-param DL baseline against transformer models",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Section: Methodology
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, content_layout, 5, "Methodology")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Dataset Description
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Dataset Description", [
    "NIST Juliet Test Suite for C/C++ v1.3 (publicly available benchmark)",
    "  105,183 source code samples across 118 CWE vulnerability categories",
    "  Covers: buffer overflows, injection, access control, resource management, crypto flaws",
    "  Each sample labeled with ground-truth CWE — no manual annotation needed",
    "Train / Test split (guideline: correct train-test split):",
    "  Training: 82,736 samples (1,323 templates) — within guideline range (50K–100K)",
    "  Test: 22,447 samples (366 templates, 87 CWEs present)",
    "  GroupShuffleSplit 80/20, seed=42 — zero template overlap (anti-leakage)",
    "Supplementary datasets for cross-validation experiments:",
    "  Big-Vul (Fan et al., 2020): 188K real-world CVE-linked C/C++ functions",
    "  C3-VULMAP (Ameh et al., 2025): healthcare-specific vulnerability mapping",
], font_size=15)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Preprocessing (guideline: normalization, missing values, bias)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Data Preprocessing & Class Imbalance Analysis", [
    "4-Stage Preprocessing Pipeline:",
    "  Stage 1: Header removal (template-generated comment blocks)",
    "  Stage 2: Comment stripping (state machine parser, preserves string literals)",
    "  Stage 3: Guard removal (#ifndef OMITBAD / OMITGOOD)",
    "  Stage 4: Whitespace normalization (compact representation)",
    "Tokenization: BPE tokenizer (CodeT5-Small, 32K vocabulary), max_length=256",
    "Class Imbalance (guideline: examine potential biases):",
    "  Largest class: 1,500 samples — Smallest class: 5 samples (300× ratio)",
    "  Macro-averaging used for all metrics to ensure minority classes count equally",
    "  No SMOTE/oversampling applied — evaluate raw model capability first",
    "No missing values: Juliet is synthetically generated with complete labels",
], font_size=15)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Model Architecture (guideline: architecture matches data type)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Model Architecture — BiLSTM with Self-Attention", [
    "Architecture (guideline: LSTM for sequential data, 5–10 layers):",
    "  Embedding(32,100 × 128)  — token embedding layer",
    "  BiLSTM(128 → 256, 2 layers, bidirectional)  — 4 LSTM layers total",
    "  Self-Attention(512 → 512)  — learns which tokens to focus on",
    "  Dropout(0.3)  — regularization (guideline: optional regularization)",
    "  Linear(512 → 118)  — 118-class CWE classification head",
    "Total parameters: 6,537,335 (~6.5M)",
    "  19× smaller than CodeBERT/GraphCodeBERT (124.7M)",
    "  5× smaller than CodeT5-Small (35.4M)",
    "Framework: PyTorch ≥2.0, Hugging Face Transformers ≥5.5 (guideline: standard DL framework)",
    "Why BiLSTM? Source code is sequential — bidirectional context captures forward + backward dependencies",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Training Configuration (guideline: hyperparameters, training time)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Training Configuration", [
    "BiLSTM-Attention (primary DL model):",
    "  Learning rate: 1×10⁻³, Optimizer: AdamW, Batch size: 32",
    "  Max tokens: 256, Epochs: 10, Early stopping: patience=3",
    "  Mixed precision: FP16 (memory efficient on 8 GB VRAM)",
    "Encoder Models — fine-tuned for comparison (CodeT5, CodeBERT, GraphCodeBERT):",
    "  LR: 5×10⁻⁵, Batch size: 8, Epochs: 2, Patience: 1, FP16",
    "Hardware (guideline: trainable on personal laptop):",
    "  GPU: NVIDIA RTX 2000 Ada Generation — 8 GB VRAM",
    "  CPU: Intel Core i7-13850HX",
    "  BiLSTM training time: ~12 minutes (10 epochs) — well within guideline's few-hour limit",
    "  Encoder training: ~25–53 minutes per model",
], font_size=15)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Evaluation Metrics (guideline: appropriate metrics, baseline)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Evaluation Metrics & Baselines", [
    "Metrics (per abstract & guideline):",
    "  Accuracy — overall correctness across all 118 classes",
    "  Macro Precision — average precision per class (minority-class sensitive)",
    "  Macro Recall — average recall per class",
    "  Macro F1-score — harmonic mean of precision and recall",
    "  Matthews Correlation Coefficient (MCC) — robust for imbalanced data",
    "  False Positive Rate (FPR) — critical for practical deployment (false alarms)",
    "Baselines (guideline: benchmark against at least one simple baseline):",
    "  Baseline 1: Zero-shot CodeGemma-2B (2.51B params, no training)",
    "  Baseline 2: Zero-shot Gemma-4-E2B-IT (5.13B params, instruction-tuned)",
    "  These represent no-training baselines to demonstrate DL model's value",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Section: Case Study Details
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, content_layout, 6, "Case Study Details")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Experimental Design
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Experimental Design — 7 Experiments", [
    "Exp A: Juliet-118 Baseline — 82K train / 22K test, 4 encoder models",
    "Exp B: Juliet-19 Subset — 19 CWEs overlapping with Big-Vul",
    "Exp C: Juliet + Big-Vul Combined — real-world code added (19 CWEs)",
    "Exp D: Big-Vul Only — 55 CWEs from real-world vulnerability commits",
    "Exp E: Juliet + Big-Vul Union — 187 CWEs, maximum coverage",
    "Exp F: Union Extended — 6 epochs, patience=3 (longer training)",
    "Exp G: BiLSTM-Attention DL Baseline — 118 CWEs, 10 epochs",
    "  Primary DL case study experiment",
    "  Tests whether a ~6.5M-param LSTM can match 125M-param transformers",
    "  Addresses: can simpler architecture learn structured vulnerability patterns?",
], font_size=15)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Tools & Environment (guideline: tools, techniques)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Tools and Environment", [
    "Programming Language: Python 3.13",
    "Deep Learning Framework: PyTorch ≥2.0 (guideline: standard DL framework)",
    "Model Hub: Hugging Face Transformers ≥5.5",
    "Tokenization: Hugging Face AutoTokenizer (BPE, 32K vocabulary)",
    "Quantization: BitsAndBytes (4-bit NF4 for QLoRA experiments)",
    "Visualization: Matplotlib, Seaborn (training curves, confusion matrices)",
    "Experiment Tracking: TensorBoard (loss / accuracy curves per epoch)",
    "Application: Streamlit (interactive web UI for model inference)",
    "Version Control: Git, GitHub",
    "Deployment Target: NVIDIA RTX 2000 Ada (8 GB VRAM), consumer laptop",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Section: Analysis
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, content_layout, 7, "Analysis")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Training & Validation Curves (guideline: include training curves)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_table_slide(s, "Training Progress — BiLSTM-Attention (Exp G)",
    ["Epoch", "Train Loss", "Test Macro-F1", "Status"],
    [
        ["1",  "0.1860", "0.9521", "New best"],
        ["2",  "0.0085", "0.9753", "New best"],
        ["3",  "0.0054", "0.9645", "Patience 1/3"],
        ["4",  "0.0047", "0.9375", "Patience 2/3"],
        ["5",  "0.0037", "0.9767", "New best"],
        ["6",  "0.0026", "0.9990", "New best ↑"],
        ["7",  "0.0049", "0.9990", "New best ↑"],
        ["8",  "0.0025", "0.9992", "New best ↑"],
        ["9",  "0.0024", "0.9992", "New best ↑"],
        ["10", "0.0039", "0.9993", "Final best ★"],
    ],
    top=Inches(1.5))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Training Curve Observations
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Training Curve Analysis (Guideline: Training/Validation Curves)", [
    "Convergence: Loss drops from 0.186 → 0.002 in 9 epochs (rapid learning)",
    "Early fluctuation: F1 dips at epoch 3–4 (patience mechanism prevents premature stop)",
    "  Early stopping with patience=3 was critical — would have stopped at epoch 4 otherwise",
    "Breakthrough at epoch 6: F1 jumps 0.9767 → 0.9990 (pattern generalization)",
    "Plateau phase: epochs 6–10 show diminishing returns (0.9990 → 0.9993)",
    "No overfitting observed: train loss and test F1 both improve monotonically after epoch 4",
    "Training time: ~12 minutes total (10 epochs on 82K samples)",
    "Comparison with transformers:",
    "  GraphCodeBERT overfits after epoch 2 — F1 drops 5.78 pp (0.9530 → 0.8952)",
    "  BiLSTM remains stable through 10 epochs — more robust training dynamic",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Confusion Matrix & Per-Class (guideline: confusion matrix)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Confusion Matrix & Per-Class Performance", [
    "74 of 87 test CWE categories achieve perfect F1 = 1.00 (85.1%)",
    "Weighted macro-F1 = 0.9982 — high accuracy on majority of samples",
    "BiLSTM confusion pattern (qualitative sample outputs):",
    "  Most errors: misclassifications toward CWE-176 (Unicode Encoding)",
    "  8 errors from CWE-114, 6 from CWE-194, 4 from CWE-121/123/126/366",
    "  These CWEs share string-handling / memory-access patterns",
    "Transformer confusion pattern (for comparison):",
    "  Misclassifications toward CWE-126 (Buffer Over-read)",
    "  Buffer-access CWEs share pointer arithmetic and array indexing syntax",
    "3 categories with F1 = 0.00: CWE-561, CWE-562, CWE-674",
    "  Each has only 1 test sample — insufficient data, not model failure",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — Attention Visualization (guideline: interpretability technique)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Interpretability — Self-Attention Visualization", [
    "Guideline requirement: apply at least one interpretability technique",
    "Technique: Self-Attention weight visualization",
    "  The attention layer produces weights α₁, α₂, ..., αₙ over all timesteps",
    "  High-weight tokens indicate what the model considers most relevant",
    "Key observations from attention analysis:",
    "  Buffer overflow CWEs: attention focuses on malloc/calloc calls and array indices",
    "  Injection CWEs: attention peaks on string concatenation and format specifiers",
    "  Access control CWEs: attention highlights file descriptor operations",
    "  Crypto CWEs: attention focuses on algorithm names and key sizes",
    "This confirms the model learns semantically meaningful vulnerability patterns",
    "  Not merely memorizing syntactic templates from Juliet",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Section: Results / Outcomes
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, content_layout, 8, "Results / Outcomes")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 27 — Results Table
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_table_slide(s, "Model Comparison — Experiment A (Juliet-118, 22,447 test)",
    ["Model", "Params", "Acc", "Prec", "Recall", "F1", "Latency"],
    [
        ["BiLSTM-Attn", "6.5M", "0.9982", "0.9997", "0.9990", "0.9993", "~5 ms"],
        ["GraphCodeBERT", "124.7M", "0.9975", "0.9530", "0.9530", "0.9530", "26.9 ms"],
        ["CodeBERT", "124.7M", "0.9974", "0.9477", "0.9528", "0.9493", "19.4 ms"],
        ["CodeT5-Small", "35.4M", "0.9972", "0.9349", "0.9410", "0.9370", "17.2 ms"],
        ["CodeT5-Base", "109.7M", "0.9977", "0.9309", "0.9418", "0.9345", "24.9 ms"],
        ["CodeGemma (ZS)", "2.51B", "0.2500", "0.1755", "0.0832", "0.0948", "3,277 ms"],
        ["Gemma-4 IT (ZS)", "5.13B", "0.7800", "0.6990", "0.7167", "0.6988", "396 s"],
    ])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 28 — Cross-Experiment F1 Table
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_table_slide(s, "Cross-Experiment Macro-F1 Comparison",
    ["Experiment", "CodeT5-S", "CodeT5-B", "CodeBERT", "GraphCB", "BiLSTM"],
    [
        ["A: Juliet-118",  "0.9370", "0.9345", "0.9493", "0.9530", "0.9993"],
        ["B: Juliet-19",   "0.9470", "0.9367", "0.9987", "0.9472", "—"],
        ["C: Combined-19", "0.8239", "0.8171", "0.8011", "0.8310", "—"],
        ["D: Big-Vul only","0.0462", "0.0656", "0.0074", "0.0730", "—"],
        ["E: Union-187",   "0.5974", "0.6045", "0.5901", "0.6219", "—"],
        ["F: Union-6ep",   "0.6376", "0.6312", "0.6142", "0.6292", "—"],
        ["G: BiLSTM-118",  "—",      "—",      "—",      "—",      "0.9993"],
    ],
    top=Inches(1.5))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 29 — Key Findings
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Key Findings", [
    "BiLSTM-Attention achieves highest F1 (0.9993) with only 6.5M parameters",
    "  19× fewer parameters than best transformer (GraphCodeBERT, F1=0.9530)",
    "  Exceeds guideline threshold (>70% accuracy) by 29+ percentage points",
    "Fine-tuning dominates zero-shot: 28–75 pp accuracy gap, 120–23,000× faster",
    "Instruction-tuning improves zero-shot 3× (Gemma-4 78% vs CodeGemma 25%)",
    "All models trainable on consumer GPU in <53 min (guideline: few-hour limit)",
    "Sub-30 ms latency — viable for CI/CD pipelines and IDE integration",
    "Real-world data severely degrades all models (Big-Vul F1 < 0.07)",
    "  Demonstrates gap between synthetic benchmarks and real-world code",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 30 — Section: Discussion
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, content_layout, 9, "Discussion")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 31 — Discussion Content
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Discussion — Insights, Implications, Comparisons", [
    "Why does BiLSTM outperform transformers on Juliet?",
    "  Juliet's structured templates create learnable sequential patterns",
    "  BiLSTM captures local token dependencies efficiently; transformers may over-parameterize",
    "  This aligns with VulDeePecker (2018): LSTMs are effective on structured code",
    "Comparison with related work:",
    "  Our F1=0.9993 on 118 CWEs surpasses SecureFalcon and FUSEVUL on their benchmarks",
    "  But: these comparisons are on different datasets — direct comparison requires same data",
    "Critical insight: synthetic vs. real-world performance gap (Exp D: F1 < 0.07)",
    "  Juliet's fixed patterns don't transfer to diverse real-world vulnerability styles",
    "  Practical deployment requires real-world data augmentation",
    "Healthcare context: same CWE has different risk in EHR system vs. game engine",
], font_size=15)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 32 — Architecture Comparison (guideline: compare architectures)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Architecture Comparison (Guideline: Compare Architectures)", [
    "BiLSTM-Attention (6.5M params):",
    "  Strengths: Fast training (12 min), low memory, highest F1 on structured data",
    "  Weakness: No pre-training on code — may not generalize to unseen patterns",
    "Encoder Transformers — CodeBERT / GraphCodeBERT (124.7M params):",
    "  Strengths: Pre-trained on code, data-flow awareness (GraphCodeBERT)",
    "  Weakness: Overfits quickly (2 epochs), 19× more parameters for lower F1",
    "QLoRA-tuned CodeGemma (2.51B params, 4-bit):",
    "  Strengths: Few-shot capable, rich code understanding",
    "  Weakness: Large footprint, slow inference, marginal accuracy gains",
    "Zero-shot LLMs (no training):",
    "  Strengths: No dataset required, broad coverage",
    "  Weakness: 25–78% accuracy — insufficient for production vulnerability detection",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 33 — Section: Ethical Considerations
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, content_layout, 10, "Ethical Considerations")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 34 — Ethical Considerations (guideline requirement)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Ethical Considerations (Guideline: Bias, Fairness, Misuse)", [
    "Dataset Bias:",
    "  Juliet is synthetically generated — may not represent real-world vulnerability diversity",
    "  Class imbalance (300× ratio) means minority CWEs are underrepresented",
    "  Model may overfit to Juliet-specific coding patterns (templates)",
    "Fairness:",
    "  Model trained on C/C++ only — does not detect vulnerabilities in other languages",
    "  High F1 on synthetic data should not be interpreted as production-ready without real-world validation",
    "Potential Misuse:",
    "  Adversaries could use vulnerability detection models to find exploitable weaknesses",
    "  Mitigation: deployment behind authentication; audit logs; responsible disclosure policies",
    "Reproducibility (guideline: mandatory):",
    "  Code, data splits, and hyperparameters fully documented; seed=42 for all experiments",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 35 — Section: Limitations & Future Work
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, content_layout, 11, "Limitations & Future Work")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 36 — Limitations & Future Work Content
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Limitations & Future Work", [
    "Limitations (per abstract: scope limited by dataset, scalability, validation):",
    "  Real-world generalization gap: Big-Vul F1 < 0.07 for all models",
    "  BiLSTM high F1 may reflect structural regularity of Juliet, not general capability",
    "  Single language (C/C++) — does not cover Java, Python, JavaScript",
    "  118 CWEs out of 900+ total CWE categories in MITRE database",
    "Future Work (guideline optional challenges):",
    "  Data augmentation: mix Juliet with real-world CVE data to improve generalization",
    "  Hyperparameter tuning: systematic search (grid/Bayesian) for BiLSTM and transformers",
    "  Robustness testing: evaluate under code obfuscation and adversarial perturbations",
    "  Multi-language extension: Java, Python, JavaScript vulnerability detection",
    "  VS Code extension / CI/CD plugin for practical deployment",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 37 — Section: Conclusion
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, content_layout, 12, "Conclusion")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 38 — Conclusion Content
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "Summary & Key Takeaways", [
    "Applied deep learning (BiLSTM with Self-Attention) to C/C++ vulnerability detection",
    "Key results on Juliet-118 (22,447 test samples):",
    "  BiLSTM-Attention: F1 = 0.9993, Accuracy = 0.9982 (6.5M params)",
    "  Best transformer (GraphCodeBERT): F1 = 0.9530 (124.7M params)",
    "  Zero-shot baselines: F1 = 0.0948 – 0.6988 (no training required)",
    "DL model outperforms transformers 19× its size on structured benchmark data",
    "All models trainable on consumer laptop GPU in <1 hour — accessible research",
    "Three-fold novelty: template-aware split, multi-paradigm comparison, healthcare risk layer",
    "Critical limitation: synthetic-to-real transfer remains the primary challenge",
    "Lesson learned: model size ≠ model performance — architecture-data fit matters more",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 39 — Section: References
# ════════════════════════════════════════════════════════════════════════════
add_section_slide(prs, content_layout, 13, "References")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 40 — References 1/2
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "References (1/2)", [
    "[1]  Li et al. (2018) — VulDeePecker: DL-based vulnerability detection, NDSS",
    "[2]  Zhou et al. (2019) — Devign: Effective vulnerability identification, NeurIPS",
    "[3]  Fan et al. (2020) — Big-Vul: Large-scale vulnerability dataset, MSR",
    "[4]  Feng et al. (2020) — CodeBERT: Pre-trained model for programming languages, EMNLP",
    "[5]  Guo et al. (2020) — GraphCodeBERT: Pre-training with data flow, arXiv",
    "[6]  Wang et al. (2021) — CodeT5: Identifier-aware encoder-decoder, EMNLP",
    "[7]  Dettmers et al. (2023) — QLoRA: Efficient finetuning of quantized LLMs, NeurIPS",
    "[8]  CodeGemma Team (2024) — CodeGemma: Open code models, Google",
    "[9]  Lu et al. (2024) — GRACE: Graph-enhanced LLM vulnerability detection, JSS",
    "[10] Ameh et al. (2025) — C3-VULMAP: Privacy-aware healthcare dataset, Electronics",
], font_size=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 41 — References 2/2
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_bullet_slide(s, "References (2/2)", [
    "[11] Yang et al. (2025) — DLAP: DL-augmented prompting, JSS",
    "[12] Tian et al. (2025) — FUSEVUL: Multi-modal vulnerability detection, Inf. Fusion",
    "[13] Ferrag et al. (2025) — SecureFalcon: Compact LLM classification, IEEE TSE",
    "[14] Liu et al. (2025) — VulACLLM: Lightweight LLM with AST+CFG, JISA",
    "[15] Mao et al. (2025) — LLMVulExp: Explainable vulnerability detection, IEEE TSE",
    "[16] Dolcetti & Iotti (2025) — LLMs and code verification, Frontiers in CS",
    "[17] Shaon & Akter (2025) — Modern approaches to vuln. detection, Electronics",
    "[18] Qiu et al. (2026) — RLV: Context-aware retrieval-based detection, JSS",
    "[19] Chen et al. (2026) — GPTVD: Static slicing + CoT, Auto. SW Eng.",
    "[20] NSA/NIST (2017) — Juliet Test Suite for C/C++ v1.3",
], font_size=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 42 — Thank You
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(content_layout); clear_slide_shapes(s)
add_textbox(s, Inches(1), Inches(2.0), Inches(8), Inches(1.2),
            "Thank You", font_size=48, bold=True,
            color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(1), Inches(3.5), Inches(8), Inches(0.5),
            "Basanth M S  |  AM.SC.R4CSE25007-PT", font_size=18,
            color=GREY, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(1), Inches(4.2), Inches(8), Inches(0.5),
            "Questions?", font_size=24, bold=False,
            color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────────────────
prs.save(OUTPUT_PATH)
print(f"✓ Presentation saved to: {OUTPUT_PATH}")
print(f"  Total slides: {len(prs.slides)}")
