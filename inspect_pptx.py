from pptx import Presentation
from pptx.util import Inches, Pt, Emu

pptx_path = r'd:\Repositories\code-vuln-detector\paper\SHARP-LLM_Secure_Lightweight_Framework_for_Source_Code_Vulnerability_Detection_and_Risk_Prioritization_in_Healthcare_Software_AM.SC.R4CSE25007-PT_v1.0.pptx'
prs = Presentation(pptx_path)

print(f'Slide width: {prs.slide_width}, height: {prs.slide_height}')
print(f'Number of slides: {len(prs.slides)}')
print(f'Number of layouts: {len(prs.slide_layouts)}')
print()

for i, layout in enumerate(prs.slide_layouts):
    print(f'Layout {i}: {layout.name}')
    for ph in layout.placeholders:
        print(f'  Placeholder idx={ph.placeholder_format.idx}, name={ph.name}, type={ph.placeholder_format.type}')
print()

for slide_idx, slide in enumerate(prs.slides):
    layout_name = slide.slide_layout.name
    print(f'--- Slide {slide_idx+1} (layout: {layout_name}) ---')
    for shape in slide.shapes:
        print(f'  Shape: {shape.shape_type}, name={shape.name}, pos=({shape.left},{shape.top}), size=({shape.width},{shape.height})')
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text[:100] if para.text else ''
                runs_info = []
                for run in para.runs:
                    font = run.font
                    ri = {
                        'text': run.text[:50],
                        'bold': font.bold,
                        'size': str(font.size) if font.size else None,
                        'name': font.name
                    }
                    try:
                        ri['color'] = str(font.color.rgb) if font.color and font.color.rgb else None
                    except:
                        ri['color'] = None
                    runs_info.append(ri)
                print(f'    Para: "{text}"')
                if runs_info:
                    print(f'    Runs: {runs_info}')
        try:
            if shape.shape_type == 13:  # PICTURE
                print(f'    [IMAGE: {shape.image.content_type}]')
        except:
            print(f'    [PICTURE shape, no embedded image]')
    print()
